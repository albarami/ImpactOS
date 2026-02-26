"""PowerPoint export — MVP-6 Section 14.

Generate deck-ready slides:
- Title slide with scenario name
- Headline numbers (GDP impact, jobs)
- Sector impact table with multipliers and direct/indirect decomposition
- Sensitivity tornado chart data
- Assumption summary
- Evidence appendix references

Uses python-pptx. Deterministic — no LLM calls.
"""

import io

from pptx import Presentation
from pptx.util import Inches, Pt


class PptxExporter:
    """Generate governed PowerPoint presentation."""

    def export(self, pack_data: dict) -> bytes:
        """Generate PPTX bytes from Decision Pack data."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._add_title_slide(prs, pack_data)
        self._add_headline_slide(prs, pack_data)
        self._add_sector_impact_slide(prs, pack_data)
        self._add_sensitivity_slide(prs, pack_data)
        self._add_assumptions_slide(prs, pack_data)
        self._add_evidence_slide(prs, pack_data)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _add_title_slide(self, prs: Presentation, data: dict) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("scenario_name", "Decision Pack")
        p.font.size = Pt(36)
        p.font.bold = True

        # Subtitle with metadata
        sub = tf.add_paragraph()
        sub.text = f"Base Year: {data.get('base_year', '')} | Currency: {data.get('currency', '')}"
        sub.font.size = Pt(18)

    def _add_headline_slide(self, prs: Presentation, data: dict) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        summary = data.get("executive_summary", {})

        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Headline Impact Numbers"
        p.font.size = Pt(28)
        p.font.bold = True

        # GDP
        gdp = summary.get("headline_gdp", 0)
        gdp_bn = gdp / 1e9
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"SAR {gdp_bn:.1f} Billion"
        p.font.size = Pt(40)
        p.font.bold = True
        sub = tf.add_paragraph()
        sub.text = "Total GDP Impact"
        sub.font.size = Pt(16)

        # Jobs
        jobs = summary.get("headline_jobs", 0)
        box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5), Inches(1.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{jobs:,} Jobs"
        p.font.size = Pt(40)
        p.font.bold = True
        sub = tf.add_paragraph()
        sub.text = "Total Employment"
        sub.font.size = Pt(16)

    def _add_sector_impact_slide(self, prs: Presentation, data: dict) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        impacts = data.get("sector_impacts", [])

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Sector Impact Breakdown"
        p.font.size = Pt(24)
        p.font.bold = True

        # Table
        rows = len(impacts) + 1  # +1 for header
        cols = 5
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5 * rows))
        table = table_shape.table

        headers = ["Sector", "Direct Impact", "Indirect Impact", "Total Impact", "Multiplier"]
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header

        for row_idx, si in enumerate(impacts, 1):
            table.cell(row_idx, 0).text = f"{si.get('sector_code', '')} - {si.get('sector_name', '')}"
            table.cell(row_idx, 1).text = f"{si.get('direct_impact', 0):,.0f}"
            table.cell(row_idx, 2).text = f"{si.get('indirect_impact', 0):,.0f}"
            table.cell(row_idx, 3).text = f"{si.get('total_impact', 0):,.0f}"
            table.cell(row_idx, 4).text = f"{si.get('multiplier', 0):.2f}"

    def _add_sensitivity_slide(self, prs: Presentation, data: dict) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sensitivity = data.get("sensitivity", [])

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Sensitivity Analysis (Tornado)"
        p.font.size = Pt(24)
        p.font.bold = True

        # Tornado data as table
        if sensitivity:
            rows = len(sensitivity) + 1
            table_shape = slide.shapes.add_table(rows, 4, Inches(1), Inches(1.5), Inches(10), Inches(0.5 * rows))
            table = table_shape.table
            headers = ["Assumption", "Low", "Base", "High"]
            for col_idx, h in enumerate(headers):
                table.cell(0, col_idx).text = h
            for row_idx, s in enumerate(sensitivity, 1):
                table.cell(row_idx, 0).text = s.get("assumption", "")
                table.cell(row_idx, 1).text = f"{s.get('low', 0):,.0f}"
                table.cell(row_idx, 2).text = f"{s.get('base', 0):,.0f}"
                table.cell(row_idx, 3).text = f"{s.get('high', 0):,.0f}"

    def _add_assumptions_slide(self, prs: Presentation, data: dict) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        assumptions = data.get("assumptions", [])

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Assumption Register"
        p.font.size = Pt(24)
        p.font.bold = True

        if assumptions:
            rows = len(assumptions) + 1
            table_shape = slide.shapes.add_table(rows, 4, Inches(1), Inches(1.5), Inches(10), Inches(0.5 * rows))
            table = table_shape.table
            headers = ["Assumption", "Value", "Range", "Status"]
            for col_idx, h in enumerate(headers):
                table.cell(0, col_idx).text = h
            for row_idx, a in enumerate(assumptions, 1):
                table.cell(row_idx, 0).text = a.get("name", "")
                table.cell(row_idx, 1).text = str(a.get("value", ""))
                rng = f"{a.get('range_min', '?')} – {a.get('range_max', '?')}"
                table.cell(row_idx, 2).text = rng
                table.cell(row_idx, 3).text = a.get("status", "")

    def _add_evidence_slide(self, prs: Presentation, data: dict) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        evidence = data.get("evidence_ledger", [])

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Evidence Appendix"
        p.font.size = Pt(24)
        p.font.bold = True

        if evidence:
            rows = len(evidence) + 1
            table_shape = slide.shapes.add_table(rows, 3, Inches(1), Inches(1.5), Inches(10), Inches(0.5 * rows))
            table = table_shape.table
            headers = ["Source", "Excerpt", "Page"]
            for col_idx, h in enumerate(headers):
                table.cell(0, col_idx).text = h
            for row_idx, e in enumerate(evidence, 1):
                table.cell(row_idx, 0).text = e.get("source", "")
                table.cell(row_idx, 1).text = e.get("text", "")
                table.cell(row_idx, 2).text = str(e.get("page", ""))
