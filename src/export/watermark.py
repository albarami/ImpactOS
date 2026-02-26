"""Watermarking service — MVP-6 Section 11.4 / 12.2.

Sandbox exports: "DRAFT — FAILS NFF GOVERNANCE" watermark on every
page/slide/sheet.
Governed exports: run_id and timestamp footer.
Watermarks are tied to the integrity signature — removing them breaks it.

Deterministic — no LLM calls.
"""

import io
from datetime import datetime
from uuid import UUID

from openpyxl import load_workbook
from openpyxl.worksheet.header_footer import HeaderFooterItem
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


SANDBOX_WATERMARK = "DRAFT \u2014 FAILS NFF GOVERNANCE"


class WatermarkService:
    """Apply watermarks to Excel and PowerPoint exports."""

    # ----- Excel -----

    def apply_sandbox_excel(self, workbook_bytes: bytes) -> bytes:
        """Apply sandbox watermark to every visible sheet header."""
        wb = load_workbook(io.BytesIO(workbook_bytes))
        for name in wb.sheetnames:
            ws = wb[name]
            if ws.sheet_state == "hidden":
                continue
            ws.oddHeader.center.text = SANDBOX_WATERMARK
            ws.oddHeader.center.size = 14
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    def apply_governed_excel(
        self,
        workbook_bytes: bytes,
        *,
        run_id: UUID,
        timestamp: datetime,
    ) -> bytes:
        """Apply governed footer with run_id and timestamp to every visible sheet."""
        wb = load_workbook(io.BytesIO(workbook_bytes))
        footer_text = f"Run: {run_id} | Generated: {timestamp.isoformat()}"
        for name in wb.sheetnames:
            ws = wb[name]
            if ws.sheet_state == "hidden":
                continue
            ws.oddFooter.center.text = footer_text
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    # ----- PowerPoint -----

    def apply_sandbox_pptx(self, pptx_bytes: bytes) -> bytes:
        """Apply sandbox watermark text box to every slide."""
        prs = Presentation(io.BytesIO(pptx_bytes))
        for slide in prs.slides:
            self._add_watermark_textbox(
                slide,
                SANDBOX_WATERMARK,
                font_size=Pt(14),
                color=RGBColor(0xFF, 0x00, 0x00),
            )
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def apply_governed_pptx(
        self,
        pptx_bytes: bytes,
        *,
        run_id: UUID,
        timestamp: datetime,
    ) -> bytes:
        """Apply governed footer text box to every slide."""
        prs = Presentation(io.BytesIO(pptx_bytes))
        footer_text = f"Run: {run_id} | {timestamp.isoformat()}"
        for slide in prs.slides:
            self._add_watermark_textbox(
                slide,
                footer_text,
                font_size=Pt(8),
                color=RGBColor(0x80, 0x80, 0x80),
                bottom=True,
            )
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    @staticmethod
    def _add_watermark_textbox(
        slide,
        text: str,
        font_size=Pt(12),
        color=RGBColor(0xFF, 0x00, 0x00),
        bottom: bool = False,
    ) -> None:
        """Add a watermark/footer text box to a slide."""
        if bottom:
            top = Inches(7.0)
        else:
            top = Inches(0.1)
        txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = True
