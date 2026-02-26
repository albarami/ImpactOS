"""Tests for PowerPoint export (MVP-6).

Covers: generate deck-ready slides with headline numbers, sector impact
data, tornado chart data, assumption summary, evidence appendix.
"""

import io

import pytest
from pptx import Presentation
from uuid_extensions import uuid7

from src.export.pptx_export import PptxExporter


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

RUN_ID = uuid7()


def _make_pack_data() -> dict:
    return {
        "run_id": str(RUN_ID),
        "scenario_name": "NEOM Logistics Zone",
        "base_year": 2023,
        "currency": "SAR",
        "executive_summary": {
            "headline_gdp": 4_200_000_000.0,
            "headline_jobs": 21200,
            "scenario_name": "NEOM Logistics Zone",
            "base_year": 2023,
            "currency": "SAR",
        },
        "sector_impacts": [
            {
                "sector_code": "C41",
                "sector_name": "Steel",
                "direct_impact": 500_000_000.0,
                "indirect_impact": 250_000_000.0,
                "total_impact": 750_000_000.0,
                "multiplier": 1.5,
            },
            {
                "sector_code": "F",
                "sector_name": "Construction",
                "direct_impact": 1_000_000_000.0,
                "indirect_impact": 800_000_000.0,
                "total_impact": 1_800_000_000.0,
                "multiplier": 1.8,
            },
        ],
        "sensitivity": [
            {"assumption": "Import Share", "low": 3.8e9, "base": 4.2e9, "high": 4.6e9},
            {"assumption": "Phasing", "low": 3.9e9, "base": 4.2e9, "high": 4.4e9},
        ],
        "assumptions": [
            {"name": "Domestic share", "value": 0.65, "range_min": 0.55, "range_max": 0.75, "status": "APPROVED"},
        ],
        "evidence_ledger": [
            {"snippet_id": str(uuid7()), "source": "SAMA Report", "text": "Steel prices +15%", "page": 42},
        ],
    }


# ===================================================================
# Presentation generation
# ===================================================================


class TestPresentationGeneration:
    """Generate PPTX presentation bytes."""

    def test_generates_bytes(self) -> None:
        exporter = PptxExporter()
        data = exporter.export(_make_pack_data())
        assert isinstance(data, bytes)
        assert len(data) > 0

    def test_valid_pptx(self) -> None:
        exporter = PptxExporter()
        data = exporter.export(_make_pack_data())
        prs = Presentation(io.BytesIO(data))
        assert len(prs.slides) >= 1


# ===================================================================
# Slide content
# ===================================================================


class TestSlideContent:
    """Deck contains required slides."""

    def test_has_title_slide(self) -> None:
        exporter = PptxExporter()
        prs = Presentation(io.BytesIO(exporter.export(_make_pack_data())))
        # First slide should be title
        slide = prs.slides[0]
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        combined = " ".join(texts)
        assert "NEOM" in combined or "Logistics" in combined

    def test_has_headline_numbers(self) -> None:
        exporter = PptxExporter()
        prs = Presentation(io.BytesIO(exporter.export(_make_pack_data())))
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text + " "
        # Should contain GDP and jobs numbers somewhere
        assert "4.2" in all_text or "4,200" in all_text or "21,200" in all_text or "21200" in all_text

    def test_has_sector_impact_table(self) -> None:
        exporter = PptxExporter()
        prs = Presentation(io.BytesIO(exporter.export(_make_pack_data())))
        # At least one slide should have a table
        has_table = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    has_table = True
                    break
        assert has_table

    def test_has_multiple_slides(self) -> None:
        exporter = PptxExporter()
        prs = Presentation(io.BytesIO(exporter.export(_make_pack_data())))
        assert len(prs.slides) >= 4

    def test_has_assumptions_slide(self) -> None:
        exporter = PptxExporter()
        prs = Presentation(io.BytesIO(exporter.export(_make_pack_data())))
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text + " "
        assert "Domestic share" in all_text or "Assumption" in all_text

    def test_has_evidence_slide(self) -> None:
        exporter = PptxExporter()
        prs = Presentation(io.BytesIO(exporter.export(_make_pack_data())))
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text + " "
        assert "Evidence" in all_text or "SAMA" in all_text
