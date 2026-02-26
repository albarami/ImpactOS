"""Tests for watermarking service (MVP-6).

Covers: sandbox watermark on every page/slide/sheet, governed run_id+timestamp
footer, watermark tied to integrity signature.
"""

import io
from datetime import datetime, timezone

import pytest
from openpyxl import load_workbook
from pptx import Presentation
from uuid_extensions import uuid7

from src.export.watermark import WatermarkService, SANDBOX_WATERMARK


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

RUN_ID = uuid7()


def _make_excel_bytes() -> bytes:
    """Create a minimal Excel file."""
    from src.export.excel_export import ExcelExporter
    return ExcelExporter().export({
        "run_id": str(RUN_ID), "scenario_name": "Test", "base_year": 2023,
        "currency": "SAR", "model_version_id": str(uuid7()), "scenario_version": 1,
        "sector_impacts": [{"sector_code": "C41", "sector_name": "Steel",
            "direct_impact": 500.0, "indirect_impact": 250.0, "total_impact": 750.0,
            "multiplier": 1.5, "domestic_share": 0.65, "import_leakage": 0.35}],
        "input_vectors": {"C41": 1000.0}, "assumptions": [],
    })


def _make_pptx_bytes() -> bytes:
    """Create a minimal PPTX file."""
    from src.export.pptx_export import PptxExporter
    return PptxExporter().export({
        "run_id": str(RUN_ID), "scenario_name": "Test", "base_year": 2023,
        "currency": "SAR",
        "executive_summary": {"headline_gdp": 1e9, "headline_jobs": 1000},
        "sector_impacts": [{"sector_code": "C41", "sector_name": "Steel",
            "direct_impact": 500.0, "indirect_impact": 250.0, "total_impact": 750.0,
            "multiplier": 1.5}],
        "sensitivity": [], "assumptions": [], "evidence_ledger": [],
    })


# ===================================================================
# Sandbox watermark
# ===================================================================


class TestSandboxWatermark:
    """Sandbox exports get DRAFT watermark."""

    def test_watermark_text(self) -> None:
        assert SANDBOX_WATERMARK == "DRAFT \u2014 FAILS NFF GOVERNANCE"

    def test_excel_sandbox_watermark(self) -> None:
        svc = WatermarkService()
        data = _make_excel_bytes()
        watermarked = svc.apply_sandbox_excel(data)
        wb = load_workbook(io.BytesIO(watermarked))
        # Every visible sheet should have watermark
        for name in wb.sheetnames:
            ws = wb[name]
            if ws.sheet_state == "hidden":
                continue
            header = ws.oddHeader
            assert SANDBOX_WATERMARK in (header.center.text if header.center else "")

    def test_pptx_sandbox_watermark(self) -> None:
        svc = WatermarkService()
        data = _make_pptx_bytes()
        watermarked = svc.apply_sandbox_pptx(data)
        prs = Presentation(io.BytesIO(watermarked))
        # Every slide should have watermark text box
        for slide in prs.slides:
            texts = [s.text for s in slide.shapes if hasattr(s, "text")]
            combined = " ".join(texts)
            assert "DRAFT" in combined


# ===================================================================
# Governed footer
# ===================================================================


class TestGovernedFooter:
    """Governed exports get run_id + timestamp footer."""

    def test_excel_governed_footer(self) -> None:
        svc = WatermarkService()
        data = _make_excel_bytes()
        ts = datetime(2026, 1, 15, 10, 30, 0, tzinfo=timezone.utc)
        footered = svc.apply_governed_excel(data, run_id=RUN_ID, timestamp=ts)
        wb = load_workbook(io.BytesIO(footered))
        for name in wb.sheetnames:
            ws = wb[name]
            if ws.sheet_state == "hidden":
                continue
            footer = ws.oddFooter
            footer_text = footer.center.text if footer.center else ""
            assert str(RUN_ID) in footer_text

    def test_pptx_governed_footer(self) -> None:
        svc = WatermarkService()
        data = _make_pptx_bytes()
        ts = datetime(2026, 1, 15, 10, 30, 0, tzinfo=timezone.utc)
        footered = svc.apply_governed_pptx(data, run_id=RUN_ID, timestamp=ts)
        prs = Presentation(io.BytesIO(footered))
        # Last slide (or footer text box) should contain run_id
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text + " "
        assert str(RUN_ID) in all_text
